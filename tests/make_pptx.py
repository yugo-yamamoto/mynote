import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
results = json.load(open(os.path.join(BASE, 'test_results.json'), encoding='utf-8'))

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

C_BG      = RGBColor(0xF9, 0xF7, 0xF4)
C_HEADER  = RGBColor(0x8B, 0x7F, 0xC0)
C_PASS    = RGBColor(0x2E, 0xA4, 0x43)
C_FAIL    = RGBColor(0xD0, 0x3B, 0x3B)
C_DARK    = RGBColor(0x2C, 0x2C, 0x2C)
C_GRAY    = RGBColor(0x88, 0x88, 0x88)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_SECTION = RGBColor(0xED, 0xEB, 0xF8)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]  # 完全ブランク

def bg(slide, color=C_BG):
    shape = slide.shapes.add_shape(1, 0, 0, SLIDE_W, SLIDE_H)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             size=18, bold=False, color=C_DARK, align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(left, top, width, height)
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb

def add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    return shape

# ============================================================
# 表紙スライド
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg(slide, RGBColor(0x8B, 0x7F, 0xC0))

add_text(slide, 'mynote', Inches(1), Inches(2.2), Inches(11), Inches(1.4),
         size=54, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
add_text(slide, 'キー操作テスト レポート',
         Inches(1), Inches(3.5), Inches(11), Inches(0.8),
         size=24, color=RGBColor(0xE0, 0xDB, 0xF5), align=PP_ALIGN.CENTER)

passed_n = sum(1 for r in results if r['pass'])
failed_n = sum(1 for r in results if not r['pass'])
add_text(slide, f'{passed_n} passed  /  {failed_n} failed  /  {len(results)} total',
         Inches(1), Inches(4.5), Inches(11), Inches(0.6),
         size=18, color=RGBColor(0xD4, 0xCF, 0xE8), align=PP_ALIGN.CENTER)

# ============================================================
# サマリスライド
# ============================================================
slide = prs.slides.add_slide(BLANK)
bg(slide)

add_rect(slide, 0, 0, SLIDE_W, Inches(1.1), C_HEADER)
add_text(slide, 'テスト結果サマリ', Inches(0.5), Inches(0.2), Inches(10), Inches(0.7),
         size=26, bold=True, color=C_WHITE)

CATEGORIES = [
    ('閲覧モード 移動',     ['01','02','03','04']),
    ('編集モード 移行',     ['05','06','07']),
    ('編集モード中のキー', ['08','09','10']),
    ('段落の削除',         ['11']),
    ('段落の移動・複製',   ['12','13','14','15']),
]

result_map = {r['id']: r for r in results}

# 利用可能高さ: 7.5 - 1.1(header) - 0.15(margin top) - 0.1(margin bottom) = 6.15 inch
# 配分: header行 0.38 + 5カテゴリ×0.26 + 15テスト行×0.27 = 0.38+1.3+4.05 = 5.73 inch → OK
CAT_H  = Inches(0.26)
ROW_H  = Inches(0.27)
HDR_H  = Inches(0.38)
col_x = [Inches(0.4), Inches(4.5), Inches(7.5), Inches(10.8)]
y = Inches(1.2)

# ヘッダ行
add_rect(slide, Inches(0.4), y, Inches(12.4), HDR_H, RGBColor(0xD8,0xD3,0xEF))
for text, x in zip(['テスト名', 'キー操作', '期待動作', '結果'],
                   [Inches(0.5), Inches(4.55), Inches(7.55), Inches(10.85)]):
    add_text(slide, text, x, y+Pt(3), Inches(3), HDR_H,
             size=10, bold=True, color=C_HEADER)
y += HDR_H

KEY_MAP = {
    '01': 'j', '02': 'k', '03': '↓', '04': '↑',
    '05': 'Enter', '06': 'Escape', '07': 'ダブルクリック',
    '08': 'Enter（編集中）', '09': 'Enter（最終段落）', '10': 'Backspace（空）',
    '11': 'Delete',
    '12': 'Alt+↓', '13': 'Alt+↑', '14': 'Shift+Alt+↓', '15': 'Shift+Alt+↑',
}

for cat_name, ids in CATEGORIES:
    # カテゴリ行
    add_rect(slide, Inches(0.4), y, Inches(12.4), CAT_H, C_SECTION)
    add_text(slide, cat_name, Inches(0.5), y+Pt(1), Inches(8), CAT_H,
             size=9, bold=True, color=C_HEADER)
    y += CAT_H
    for id_ in ids:
        r = result_map.get(id_)
        if not r: continue
        c = C_PASS if r['pass'] else C_FAIL
        badge = '✓ PASS' if r['pass'] else '✗ FAIL'
        add_text(slide, r['name'], col_x[0]+Inches(0.1), y+Pt(2), Inches(3.8), ROW_H, size=9, color=C_DARK)
        add_text(slide, KEY_MAP.get(id_,''), col_x[1]+Inches(0.1), y+Pt(2), Inches(2.8), ROW_H, size=9, color=C_DARK)
        desc_short = r['description'][:36] + ('…' if len(r['description']) > 36 else '')
        add_text(slide, desc_short, col_x[2]+Inches(0.1), y+Pt(2), Inches(3.1), ROW_H, size=8, color=C_GRAY)
        add_rect(slide, col_x[3]+Inches(0.05), y+Pt(3), Inches(1.2), ROW_H-Pt(4), c)
        add_text(slide, badge, col_x[3]+Inches(0.08), y+Pt(4), Inches(1.1), ROW_H,
                 size=9, bold=True, color=C_WHITE)
        y += ROW_H

# ============================================================
# テストケース個別スライド (スクリーンショット付き)
# ============================================================
SECTION_TITLES = {
    '01': '閲覧モード 移動', '02': '閲覧モード 移動',
    '03': '閲覧モード 移動', '04': '閲覧モード 移動',
    '05': '編集モード 移行', '06': '編集モード 移行', '07': '編集モード 移行',
    '08': '編集モード中のキー', '09': '編集モード中のキー', '10': '編集モード中のキー',
    '11': '段落の削除',
    '12': '段落の移動・複製', '13': '段落の移動・複製',
    '14': '段落の移動・複製', '15': '段落の移動・複製',
}

for r in results:
    slide = prs.slides.add_slide(BLANK)
    bg(slide)

    # ヘッダ帯
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.0), C_HEADER)
    section = SECTION_TITLES.get(r['id'], '')
    add_text(slide, f"#{r['id']}  {section}", Inches(0.4), Inches(0.05),
             Inches(9), Inches(0.4), size=13, color=RGBColor(0xD4,0xCF,0xE8))
    add_text(slide, r['name'], Inches(0.4), Inches(0.42),
             Inches(10), Inches(0.55), size=22, bold=True, color=C_WHITE)

    # PASS/FAIL バッジ
    badge_c = C_PASS if r['pass'] else C_FAIL
    add_rect(slide, Inches(11.2), Inches(0.28), Inches(1.7), Inches(0.48), badge_c)
    add_text(slide, '✓  PASS' if r['pass'] else '✗  FAIL',
             Inches(11.25), Inches(0.32), Inches(1.6), Inches(0.44),
             size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    # スクリーンショット (左側)
    img_path = r.get('screenshotPath', '')
    if img_path and os.path.exists(img_path):
        slide.shapes.add_picture(img_path,
            Inches(0.35), Inches(1.15), Inches(7.8), Inches(5.25))

    # 説明パネル (右側)
    add_rect(slide, Inches(8.35), Inches(1.15), Inches(4.6), Inches(5.25),
             RGBColor(0xFF,0xFF,0xFF), RGBColor(0xD8,0xD3,0xEF))

    add_text(slide, '操作キー', Inches(8.55), Inches(1.3), Inches(4.0), Inches(0.35),
             size=11, bold=True, color=C_HEADER)
    add_rect(slide, Inches(8.55), Inches(1.65), Inches(4.0), Inches(0.52),
             RGBColor(0xED,0xEB,0xF8))
    add_text(slide, KEY_MAP.get(r['id'], '—'),
             Inches(8.65), Inches(1.68), Inches(3.8), Inches(0.48),
             size=16, bold=True, color=C_DARK)

    add_text(slide, 'テストの説明', Inches(8.55), Inches(2.32), Inches(4.0), Inches(0.35),
             size=11, bold=True, color=C_HEADER)
    add_text(slide, r['description'],
             Inches(8.55), Inches(2.68), Inches(4.0), Inches(2.0),
             size=12, color=C_DARK, wrap=True)

    add_text(slide, '期待値', Inches(8.55), Inches(4.75), Inches(4.0), Inches(0.35),
             size=11, bold=True, color=C_HEADER)
    expected_texts = {
        '01': '次段落 (index=1) にフォーカスが移動する',
        '02': '前段落 (index=0) にフォーカスが戻る',
        '03': '次段落 (index=1) にフォーカスが移動する',
        '04': '前段落 (index=0) にフォーカスが戻る',
        '05': 'contenteditable = true になる',
        '06': 'contenteditable = false になる',
        '07': 'contenteditable = true になる',
        '08': 'フォーカスが index=1 の段落に移動する',
        '09': '段落数が 1 増加する',
        '10': '段落数が 1 減少する',
        '11': '段落数が 1 減少する',
        '12': '段落の順序が入れ替わる（下に移動）',
        '13': '段落の順序が入れ替わる（上に移動）',
        '14': '段落数が 1 増加し、直下に同じテキストが複製される',
        '15': '段落数が 1 増加し、直上に同じテキストが複製される',
    }
    add_text(slide, expected_texts.get(r['id'], '—'),
             Inches(8.55), Inches(5.1), Inches(4.0), Inches(0.8),
             size=11, color=C_DARK, wrap=True)

    if not r['pass'] and r.get('error'):
        add_text(slide, f"エラー: {r['error']}",
                 Inches(8.55), Inches(5.9), Inches(4.0), Inches(0.45),
                 size=10, color=C_FAIL)

OUT = os.path.join(BASE, 'mynote_test_report.pptx')
prs.save(OUT)
print(f'保存完了: {OUT}')
